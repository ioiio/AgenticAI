"""
Generates a highly professional widescreen (16:9) PowerPoint presentation (.pptx)
for a 45-minute lecture on 'How LLMs Work Underneath'.

Requirements:
    pip install python-pptx

Author: Gemini 3 Flash Paid Tier
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -----------------------------------------------------------------------------
# Color Palette Definition (Premium Dark Tech Theme)
# -----------------------------------------------------------------------------
BG_COLOR = RGBColor(11, 15, 25)         # Deep Navy Black
CARD_BG = RGBColor(22, 27, 34)          # Card Dark Gray
ACCENT_BLUE = RGBColor(88, 166, 255)    # Electric Blue
ACCENT_MINT = RGBColor(86, 227, 159)    # Mint Green
TEXT_WHITE = RGBColor(240, 246, 252)    # Crisp White
TEXT_MUTED = RGBColor(139, 148, 158)    # Light Gray
HIGHLIGHT_RED = RGBColor(249, 130, 108)  # Coral Red

# -----------------------------------------------------------------------------
# Helper Functions for Visual Elements
# -----------------------------------------------------------------------------
def set_slide_background(slide, color):
    """Overrides slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=ACCENT_BLUE):
    """Creates a styled card container shape acting as a clean content tile."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def add_title(slide, text):
    """Standardizes slide title positioning, fonts, and boundaries."""
    tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.13), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TEXT_WHITE
    return tx_box

def add_text_box(slide, text, left, top, width, height, font_size=16, bold=False, color=TEXT_WHITE, align=PP_ALIGN.LEFT):
    """Convenience helper to output wrapped, styled text onto slides."""
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx_box

def add_bullet_list(slide, items, left, top, width, height, item_size=16):
    """Draws custom clean bullet points inside a defined text bounds."""
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run_bullet = p.add_run()
        run_bullet.text = "▪ "
        run_bullet.font.name = "Segoe UI"
        run_bullet.font.size = Pt(item_size)
        run_bullet.font.bold = True
        run_bullet.font.color.rgb = ACCENT_MINT
        
        run_text = p.add_run()
        run_text.text = item
        run_text.font.name = "Segoe UI"
        run_text.font.size = Pt(item_size)
        run_text.font.color.rgb = TEXT_WHITE

# -----------------------------------------------------------------------------
# Core Presentation Orchestration
# -----------------------------------------------------------------------------
def main():
    prs = Presentation()
    
    # Force standard 16:9 widescreen slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layout (index 6 is typically completely blank)
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------------------
    # Slide 1: Premium Title Slide
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_COLOR)
    
    # Left decorative accent bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    
    # Topic Title
    tx_title = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.0))
    tf = tx_title.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "How LLMs Work Underneath\n"
    r1.font.name = "Segoe UI"
    r1.font.size = Pt(54)
    r1.font.bold = True
    r1.font.color.rgb = TEXT_WHITE
    
    r2 = p1.add_run()
    r2.text = "Tokens, Embeddings, Context, and the Stateless API Paradigm"
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(22)
    r2.font.color.rgb = ACCENT_MINT
    
    # Metadata Subtitle
    add_text_box(slide1, 
                 "Course: AI Engineering  |  Duration: 45-Minute Deep Dive Lecture\nInstructor: Prakash Tripathi", 
                 Inches(1.2), Inches(4.5), Inches(10.0), Inches(1.0), 
                 font_size=15, color=TEXT_MUTED)

    # -------------------------------------------------------------------------
    # Slide 2: Lecture Agenda & Timings (45 Mins Total)
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_COLOR)
    add_title(slide2, "Session Structure & Core Timings")
    
    add_text_box(slide2, "A tactical timeline designed to master operational API realities.", 
                 Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.4), font_size=16, color=TEXT_MUTED)
    
    # 5 Sequential Agenda Steps with accurate times
    agenda_data = [
        ("M1: Tokenization (10 Mins)", "Sub-word splitting, Byte-Pair Encoding (BPE), and vocabulary metrics."),
        ("M2: Vector Embeddings (10 Mins)", "Mapping natural language coordinates & computing Cosine Proximity."),
        ("M3: Context Windows (10 Mins)", "Quadratic limits, Attention cost scaling, and RAM constraints."),
        ("M4: Sampling & Controls (10 Mins)", "Temperature smoothing distribution bounds and Top-P cutoff logic."),
        ("M5: Stateless Paradigms (5 Mins)", "Converting dialogue contexts and contrasting raw completed strings.")
    ]
    
    start_left = Inches(0.6)
    width = Inches(2.3)
    gap = Inches(0.15)
    
    for idx, (title, desc) in enumerate(agenda_data):
        tile_left = start_left + idx * (width + gap)
        create_card(slide2, tile_left, Inches(1.8), width, Inches(4.8), CARD_BG, ACCENT_BLUE)
        
        # Draw small step icon tag
        tag = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, tile_left + Inches(0.2), Inches(2.0), Inches(1.9), Inches(0.4))
        tag.fill.solid()
        tag.fill.fore_color.rgb = ACCENT_BLUE
        tag.line.fill.background()
        tf_tag = tag.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.alignment = PP_ALIGN.CENTER
        r_tag = p_tag.add_run()
        r_tag.text = f"Step {idx+1}"
        r_tag.font.name = "Segoe UI"
        r_tag.font.size = Pt(14)
        r_tag.font.bold = True
        r_tag.font.color.rgb = BG_COLOR
        
        # Add Module Title
        add_text_box(slide2, title, tile_left + Inches(0.15), Inches(2.6), width - Inches(0.3), Inches(1.0), 
                     font_size=15, bold=True, color=ACCENT_MINT)
        
        # Add Module Details
        add_text_box(slide2, desc, tile_left + Inches(0.15), Inches(3.6), width - Inches(0.3), Inches(2.8), 
                     font_size=13, color=TEXT_WHITE)

    # -------------------------------------------------------------------------
    # Slide 3: Module 1 - Tokenization
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_COLOR)
    add_title(slide3, "Understanding Tokenization")
    
    # Split into 2 primary column cards
    create_card(slide3, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide3, "The Sub-Word Translation Layer", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_tok_left = [
        "LLMs don't read words directly; raw characters are translated into sub-word token chunks.",
        "A typical rule-of-thumb: 1 token roughly maps to 4 characters (or 0.75 English words).",
        "Tokens prevent spelling vulnerabilities while efficiently compressing compound strings."
    ]
    add_bullet_list(slide3, bullets_tok_left, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    create_card(slide3, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide3, "How BPE Constructs Vocabularies", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    bullets_tok_right = [
        "Byte-Pair Encoding (BPE) iterative algorithms find recurring byte strings to index into a static dictionary.",
        "Vocabulary metrics usually define dictionary bounds (e.g., 32,000 up to 100,000+ distinct words).",
        "Code expressions, whitespace, and variable sequences utilize distinct mathematical token blocks."
    ]
    add_bullet_list(slide3, bullets_tok_right, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)

    # -------------------------------------------------------------------------
    # Slide 4: Module 2 - Vector Embeddings
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, BG_COLOR)
    add_title(slide4, "Embedding Models & Spatial Math")
    
    # Left Card: Semantic Projections
    create_card(slide4, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide4, "Multidimensional Vector Maps", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_emb = [
        "Embeddings translate static tokens into mathematical coordinates inside dense vector spaces (e.g., 1536 dims).",
        "Similar concepts are mathematically close together, linking 'puppy' and 'dog' spatially.",
        "These vector spaces are constructed automatically during the model's massive baseline pre-training runs."
    ]
    add_bullet_list(slide4, bullets_emb, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    # Right Card: Cosine Proximity Formula Focus
    create_card(slide4, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_MINT)
    add_text_box(slide4, "Measuring Semantic Distance", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    
    formula_text = (
        "Semantic similarity is computed mathematically via spatial trigonometry:\n\n"
        "             A  ·  B\n"
        "   cos(θ) = -----------\n"
        "            ||A|| ||B||\n\n"
        "• Range is [-1, 1], where 1 denotes identical direction vectors.\n"
        "• Vector search databases index these distances to enable superfast retrieval (RAG)."
    )
    add_text_box(slide4, formula_text, Inches(7.2), Inches(2.4), Inches(5.2), Inches(3.8), font_size=16, color=TEXT_WHITE)

    # -------------------------------------------------------------------------
    # Slide 5: Module 3 - Context Windows & Attention Limits
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, BG_COLOR)
    add_title(slide5, "Context Windows & Scaling Laws")
    
    create_card(slide5, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide5, "The Ephemeral Memory Limit", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_ctx = [
        "The context window determines the absolute total threshold of inputs and output tokens the LLM can process in one pass.",
        "Exceeding this rigid window boundary triggers immediate truncation, resulting in chat state memory loss.",
        "While modern models support 128K+ windows, older systems were constrained to 4K or 8K."
    ]
    add_bullet_list(slide5, bullets_ctx, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    create_card(slide5, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide5, "Quadratic Computational Cost", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=HIGHLIGHT_RED)
    
    cost_text = (
        "Standard Transformer self-attention costs scale quadratically with token length:\n\n"
        "            O( N² ) Complexity\n\n"
        "• Doubling input sequence length quadruples the key-value cache memory required on GPUs.\n"
        "• This computational bottleneck is why infinite context is physically and financially impossible, making efficient context management a core engineering skill."
    )
    add_text_box(slide5, cost_text, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8), font_size=16, color=TEXT_WHITE)

    # -------------------------------------------------------------------------
    # Slide 6: Module 4 - Next-Token Decoding
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, BG_COLOR)
    add_title(slide6, "Decoding: Next-Token Selection")
    
    create_card(slide6, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide6, "The Probabilistic Choice Mechanism", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_dec = [
        "LLMs don't write blocks of text; they generate content step-by-step by predicting the single next best token.",
        "After processing the context, the model outputs raw scores (logits) across its entire vocabulary.",
        "A softmax function maps these raw logits to a beautiful 0% to 100% probability distribution.",
        "The decoding strategy determines how we sample from this list of candidates."
    ]
    add_bullet_list(slide6, bullets_dec, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    create_card(slide6, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide6, "Sampling Hyperparameters", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    bullets_dec_right = [
        "Sampling parameters act as controls on this vocabulary selection pool.",
        "They allow developers to balance cold, deterministic code execution rules with open, creative reasoning paths.",
        "Configuring these options is critical to building robust production pipelines."
    ]
    add_bullet_list(slide6, bullets_dec_right, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)

    # -------------------------------------------------------------------------
    # Slide 7: Sampling Visualized - Temperature
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, BG_COLOR)
    add_title(slide7, "Sampling Parameters: Temperature")
    
    add_text_box(slide7, "Controlling output distribution curves directly.", 
                 Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.4), font_size=16, color=TEXT_MUTED)
    
    # 3 Column comparative cards
    temp_data = [
        ("Deterministic (T = 0.1)", "Highly predictable, repeatable results.", "The probability of the absolute top token spikes dramatically while the tail is suppressed.\n\nIdeal for code execution, schema constraints, and technical formatting."),
        ("Balanced (T = 0.7)", "Default conversational baseline.", "Maintains soft variations, mimicking natural human cadence.\n\nStandard baseline for email drafting, general chat interfaces, and analysis."),
        ("High-Entropy (T = 1.2)", "Uninhibited creative exploration.", "Flattens the distribution curve, elevating unlikely token candidates.\n\nGreat for brainstorming or creative writing; highly prone to hallucinations.")
    ]
    
    width_col = Inches(3.8)
    gap_col = Inches(0.3)
    
    for idx, (head, sub, details) in enumerate(temp_data):
        col_left = Inches(0.6) + idx * (width_col + gap_col)
        create_card(slide7, col_left, Inches(1.8), width_col, Inches(4.8), CARD_BG, ACCENT_BLUE)
        
        add_text_box(slide7, head, col_left + Inches(0.2), Inches(2.1), width_col - Inches(0.4), Inches(0.4), font_size=18, bold=True, color=ACCENT_MINT)
        add_text_box(slide7, sub, col_left + Inches(0.2), Inches(2.6), width_col - Inches(0.4), Inches(0.4), font_size=14, bold=True, color=TEXT_MUTED)
        add_text_box(slide7, details, col_left + Inches(0.2), Inches(3.2), width_col - Inches(0.4), Inches(3.2), font_size=14, color=TEXT_WHITE)

    # -------------------------------------------------------------------------
    # Slide 8: Sampling Parameters - Top-P
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, BG_COLOR)
    add_title(slide8, "Sampling Parameters: Top-P")
    
    create_card(slide8, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide8, "Nucleus Boundary Constraints", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_topp = [
        "Top-P (Nucleus Sampling) restricts the selection pool based on cumulative probability percentage.",
        "A Top-P of 0.90 isolates only the high-probability tokens that make up the top 90% of the distribution.",
        "The model dynamically expands or shrinks the candidate pool based on its confidence score.",
        "Unlike temperature, Top-P completely clips off the long, cold tail of logical nonsense."
    ]
    add_bullet_list(slide8, bullets_topp, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    create_card(slide8, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide8, "Developer Tuning Guidance", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    bullets_topp_right = [
        "Strict Rule: Never adjust both Temperature and Top-P concurrently. Lock one down first.",
        "To enforce rigid constraints (like JSON output), set Temperature to 0.0 or Top-P to 0.05.",
        "For expressive writing benchmarks, hold Temperature at 0.8 and scale Top-P to 0.95."
    ]
    add_bullet_list(slide8, bullets_topp_right, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)

    # -------------------------------------------------------------------------
    # Slide 9: Module 5 - Stateless Architecture
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, BG_COLOR)
    add_title(slide9, "The Stateless API Reality")
    
    # Large focal card on the left
    create_card(slide9, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide9, "Zero Stored Session State", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_state = [
        "LLM Chat APIs are completely stateless. The remote inference engine retains zero memory of previous API calls.",
        "Every single call is processed independently, with no context carried over natively.",
        "This architecture allows standard serverless compute infrastructure to scale massively."
    ]
    add_bullet_list(slide9, bullets_state, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    # Visual simulation card on the right
    create_card(slide9, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_MINT)
    add_text_box(slide9, "How Chat History is Built", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    
    simulation_text = (
        "To simulate ongoing conversational context:\n\n"
        "1. The application must store message history database records locally.\n"
        "2. For every new message, compile the entire chat history into an array.\n"
        "3. Send the full array back to the stateless API.\n\n"
        "This means API payloads grow larger with each conversational turn."
    )
    add_text_box(slide9, simulation_text, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8), font_size=16, color=TEXT_WHITE)

    # -------------------------------------------------------------------------
    # Slide 10: Legacy vs. Chat Completions
    # -------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, BG_COLOR)
    add_title(slide10, "Legacy vs. Chat Completions")
    
    add_text_box(slide10, "Contrasting legacy endpoints with modern structural JSON role-based array APIs.", 
                 Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.4), font_size=16, color=TEXT_MUTED)
    
    # 2 Column cards comparing schemas
    create_card(slide10, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE)
    add_text_box(slide10, "Legacy: /v1/completions", Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    legacy_text = (
        "Payload: Single raw text prompt string\n\n"
        "• Lacked structure, requiring manual tags (e.g., 'User: / AI:') in the prompt.\n"
        "• Highly vulnerable to context pollution and prompt injections.\n"
        "• Difficult to enforce clear system-level instructions or guardrails."
    )
    add_text_box(slide10, legacy_text, Inches(0.9), Inches(2.8), Inches(5.2), Inches(3.4), font_size=15, color=TEXT_WHITE)
    
    create_card(slide10, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE)
    add_text_box(slide10, "Modern: /v1/chat/completions", Inches(7.2), Inches(2.1), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    chat_text = (
        "Payload: Structured Array of Role-based Objects\n\n"
        "• Separation of Roles: 'system' instructions, 'user' inputs, and 'assistant' replies.\n"
        "• System prompts provide solid, injection-resistant guardrails.\n"
        "• Enables native JSON validation schemas (e.g., Pydantic parsing)."
    )
    add_text_box(slide10, chat_text, Inches(7.2), Inches(2.8), Inches(5.2), Inches(3.4), font_size=15, color=TEXT_WHITE)

    # -------------------------------------------------------------------------
    # Slide 11: Developer Takeaways & Practical Synthesis
    # -------------------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, BG_COLOR)
    add_title(slide11, "Best Practices & Synthesis")
    
    # Split summary list
    create_card(slide11, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide11, "Cost & Performance Optimization", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_BLUE)
    bullets_dev_l = [
        "Be mindul of token counts: Large few-shot prompts consume substantial context window bandwidth.",
        "Implement aggressive message pruning or summaries to prevent context overhead costs.",
        "Use local tokenizer tools (like Tiktoken) to monitor payloads before hitting APIs."
    ]
    add_bullet_list(slide11, bullets_dev_l, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)
    
    create_card(slide11, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), CARD_BG, ACCENT_BLUE)
    add_text_box(slide11, "Reliability & Safety Tuning", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), font_size=20, bold=True, color=ACCENT_MINT)
    bullets_dev_r = [
        "Enforce strict schema validation on structural APIs using Pydantic outputs.",
        "Lock Temperature down to 0.0 for factual production databases.",
        "Always isolate sensitive system instructions from untrusted user inputs to prevent injection exploits."
    ]
    add_bullet_list(slide11, bullets_dev_r, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8), item_size=15)

    # -------------------------------------------------------------------------
    # Slide 12: Interactive Q&A Session (15 Mins)
    # -------------------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, BG_COLOR)
    
    # Centered large content cards
    box = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.73), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = ACCENT_BLUE
    box.line.width = Pt(2.0)
    
    add_text_box(slide12, "Interactive Discussion & Q&A", Inches(1.0), Inches(2.3), Inches(11.33), Inches(0.8), font_size=40, bold=True, color=ACCENT_MINT, align=PP_ALIGN.CENTER)
    
    details_qa = (
        "Time Allocated: 15 Minutes\n\n"
        "• Troubleshoot tokenization limits and rate limits.\n"
        "• Discuss cost optimization strategies for multi-turn chats.\n"
        "• Deep-dive into Temperature vs. Top-P behaviors."
    )
    add_text_box(slide12, details_qa, Inches(1.0), Inches(3.3), Inches(11.33), Inches(2.0), font_size=18, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide12, "AI Engineering Course | Instructor: Prakash Tripathi", Inches(0.8), Inches(5.8), Inches(11.73), Inches(0.5), font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------------
    # Save Presentation
    # -------------------------------------------------------------------------
    output_filename = "LLMs_Under_The_Hood_45Min.pptx"
    prs.save(output_filename)
    print(f"\n[Success] Your 45-minute slide deck was generated: '{output_filename}'")

if __name__ == "__main__":
    main()